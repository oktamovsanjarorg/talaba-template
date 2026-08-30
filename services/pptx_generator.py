import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def create_presentation_pptx(slides_data: list, topic: str, output_path: str, student_name: str = "Talaba"):
    """
    16:9 formatdagi ko'p ustunli (Cards Grid), zamonaviy tipografiyali va premium dizayndagi PowerPoint generatori.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    # Designer ranglar palitrasi
    COLOR_DARK_BG = RGBColor(11, 19, 43)        # #0B132B Dark Navy
    COLOR_PRIMARY = RGBColor(30, 58, 138)       # #1E3A8A Deep Blue
    COLOR_ACCENT_BLUE = RGBColor(37, 99, 235)   # #2563EB Electric Blue
    COLOR_ACCENT_TEAL = RGBColor(13, 148, 136)  # #0D9488 Emerald/Teal
    COLOR_ACCENT_VIOLET = RGBColor(124, 58, 237)# #7C3AED Violet
    COLOR_ACCENT_AMBER = RGBColor(217, 119, 6)  # #D97706 Amber
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_LIGHT_BG = RGBColor(248, 250, 252)    # #F8FAFC Slate Light
    COLOR_CARD_BG = RGBColor(255, 255, 255)     # Oq kartochka
    COLOR_CARD_BORDER = RGBColor(226, 232, 240) # #E2E8F0
    COLOR_TEXT_MAIN = RGBColor(15, 23, 42)      # #0F172A Slate 900
    COLOR_TEXT_MUTED = RGBColor(71, 85, 105)    # #475569 Slate 600

    ACCENT_COLORS = [COLOR_ACCENT_BLUE, COLOR_ACCENT_TEAL, COLOR_ACCENT_VIOLET, COLOR_ACCENT_AMBER]
    blank_layout = prs.slide_layouts[6]

    # ================= 1. TITUL SLAYDI (EXECUTIVE DARK) =================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Orqa fon
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK_BG
    bg1.line.fill.background()

    # Yuqori nishon (Badge)
    top_badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(3.2), Inches(0.45))
    top_badge.fill.solid()
    top_badge.fill.fore_color.rgb = COLOR_ACCENT_BLUE
    top_badge.line.fill.background()
    tf_tb = top_badge.text_frame
    p_tb = tf_tb.paragraphs[0]
    p_tb.text = "AKADEMIK TAQDIMOT | 2026"
    p_tb.font.size = Pt(13)
    p_tb.font.bold = True
    p_tb.font.color.rgb = COLOR_WHITE
    p_tb.alignment = PP_ALIGN.CENTER

    # Sarlavha qutisi
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = topic.upper()
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.font.name = "Arial"

    # Muallif paneli
    author_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.2), Inches(6.5), Inches(1.2))
    author_card.fill.solid()
    author_card.fill.fore_color.rgb = RGBColor(26, 38, 70)
    author_card.line.color.rgb = COLOR_ACCENT_BLUE
    author_card.line.width = Pt(1.5)
    tf_ac = author_card.text_frame
    p_ac1 = tf_ac.paragraphs[0]
    p_ac1.text = f"👤 Tayyorladi: {student_name}"
    p_ac1.font.size = Pt(17)
    p_ac1.font.bold = True
    p_ac1.font.color.rgb = COLOR_WHITE
    p_ac2 = tf_ac.add_paragraph()
    p_ac2.text = "🏛 Oliy Ta'lim Muassasasi | Ilmiy Seminar Taqdimoti"
    p_ac2.font.size = Pt(13)
    p_ac2.font.color.rgb = RGBColor(148, 163, 184)

    # ================= 2. MAZMUN SLAYDLARI (CARD GRID) =================
    total_content_slides = len(slides_data)
    for idx, slide_item in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)

        # Orqa fon (Soft Platinum)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_LIGHT_BG
        bg.line.fill.background()

        # Yuqori Header paneli
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.3))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = COLOR_PRIMARY
        header_bar.line.fill.background()

        # Raqam pill (01, 02)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.35), Inches(0.9), Inches(0.6))
        pill.fill.solid()
        pill.fill.fore_color.rgb = ACCENT_COLORS[(idx - 1) % len(ACCENT_COLORS)]
        pill.line.fill.background()
        tf_p = pill.text_frame
        p_p = tf_p.paragraphs[0]
        p_p.text = f"{idx:02d}"
        p_p.font.size = Pt(18)
        p_p.font.bold = True
        p_p.font.color.rgb = COLOR_WHITE
        p_p.alignment = PP_ALIGN.CENTER

        # Sarlavha
        title_box = slide.shapes.add_textbox(Inches(1.9), Inches(0.25), Inches(10.5), Inches(0.8))
        tf_t = title_box.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = slide_item.get("title", f"{idx}-Mavzu")
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.font.name = "Arial"

        # Subtitle / Key Takeaway
        subtitle_text = slide_item.get("subtitle") or slide_item.get("summary") or f"{topic} bo'yicha tahliliy ko'rsatkichlar va amaliy xulosalar"
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.5))
        tf_s = sub_box.text_frame
        p_s = tf_s.paragraphs[0]
        p_s.text = f"💡 {subtitle_text}"
        p_s.font.size = Pt(14)
        p_s.font.italic = True
        p_s.font.color.rgb = COLOR_TEXT_MUTED
        p_s.font.name = "Arial"

        # KARTALAR (Cards Grid Layout)
        cards = slide_item.get("cards", [])
        if not cards and "bullets" in slide_item:
            # Fallback agar faqat bullets bo'lsa
            cards = [{"card_title": f"Fakt {i+1}", "card_text": b} for i, b in enumerate(slide_item["bullets"])]

        num_cards = min(len(cards), 4)
        if num_cards == 0:
            num_cards = 1
            cards = [{"card_title": "Asosiy Mazmun", "card_text": "Taqdimot bo'yicha ma'lumotlar o'rganilmoqda."}]

        # 3 ta ustunli yoki 4 ta ustunli grid hisoblash
        total_width = 11.7
        gap = 0.25
        card_w = (total_width - (num_cards - 1) * gap) / num_cards
        card_h = 4.4
        top_pos = 2.05

        for c_idx in range(num_cards):
            c_data = cards[c_idx]
            left_pos = 0.8 + c_idx * (card_w + gap)
            accent_color = ACCENT_COLORS[c_idx % len(ACCENT_COLORS)]

            # Oq kartochka foni
            card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(top_pos), Inches(card_w), Inches(card_h))
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = COLOR_CARD_BG
            card_shape.line.color.rgb = COLOR_CARD_BORDER
            card_shape.line.width = Pt(1.2)

            # Kartochka yuqori aksent chizig'i
            top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left_pos), Inches(top_pos), Inches(card_w), Inches(0.12))
            top_line.fill.solid()
            top_line.fill.fore_color.rgb = accent_color
            top_line.line.fill.background()

            # Kartochka ichidagi matn
            c_box = slide.shapes.add_textbox(Inches(left_pos + 0.15), Inches(top_pos + 0.3), Inches(card_w - 0.3), Inches(card_h - 0.5))
            tf_c = c_box.text_frame
            tf_c.word_wrap = True

            # Kartochka sarlavhasi
            p_ct = tf_c.paragraphs[0]
            p_ct.text = c_data.get("card_title") or f"Xususiyat {c_idx+1}"
            p_ct.font.size = Pt(16)
            p_ct.font.bold = True
            p_ct.font.color.rgb = COLOR_TEXT_MAIN
            p_ct.font.name = "Arial"

            # Kartochka tavsifi
            p_cd = tf_c.add_paragraph()
            p_cd.text = f"\n{c_data.get('card_text') or ''}"
            p_cd.font.size = Pt(13)
            p_cd.font.color.rgb = COLOR_TEXT_MUTED
            p_cd.font.name = "Arial"

        # Footer (Sahifa raqami)
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
        tf_f = footer_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"Talaba AI Presentation System  |  {topic[:40]}  |  Slayd {idx + 1} / {total_content_slides + 2}"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = RGBColor(148, 163, 184)
        p_f.alignment = PP_ALIGN.RIGHT

    # ================= 3. YAKUNIY SLAYD =================
    slide_end = prs.slides.add_slide(blank_layout)
    bg_end = slide_end.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg_end.fill.solid()
    bg_end.fill.fore_color.rgb = COLOR_DARK_BG
    bg_end.line.fill.background()

    end_box = slide_end.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.0))
    tf_end = end_box.text_frame
    p_end = tf_end.paragraphs[0]
    p_end.alignment = PP_ALIGN.CENTER
    p_end.text = "E'TIBORINGIZ UCHUN RAHMAT!"
    p_end.font.size = Pt(40)
    p_end.font.bold = True
    p_end.font.color.rgb = COLOR_WHITE
    p_end.font.name = "Arial"

    p_sub = tf_end.add_paragraph()
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "\nSavollar va ilmiy mulohazalar uchun minnatdormiz"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_ACCENT_TEAL
    p_sub.font.name = "Arial"

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
